import os
import io
import streamlit as st
from PyPDF2 import PdfReader
from youtube_transcript_api import YouTubeTranscriptApi
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from pptx import Presentation

load_dotenv()

# Set the API key
GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')
genai.configure(api_key=GOOGLE_API_KEY)

def get_pdf_text(pdf_bytes):
    """Extracts text from uploaded PDFs."""
    text = ""
    pdf_stream = io.BytesIO(pdf_bytes)
    pdf_reader = PdfReader(pdf_stream)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_transcript_from_captions(url):
    """Retrieves transcript from YouTube video captions."""
    try:
        video_id = url.split('=')[1]
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = ""
        for line in transcript:
            text += " " + line["text"] + "\n"
        return text
    except Exception as e:
        st.error(f"Error retrieving transcript: {str(e)}")
        return None


def process_google_slides(slide_data):
    """Processes Google Slide data."""
    try:
        # Read the presentation
        prs = Presentation(io.BytesIO(slide_data.read()))

        # Extract text from each slide
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"Error processing Google Slide data: {str(e)}")
        return None

def get_text_chunks(text):
    """Splits text into chunks for vectorization."""
    if not text:
        return []
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    """Creates or loads the FAISS vector store."""
    #clear_vector_store()  # Clear the vector store
    if not text_chunks:
        return None
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=GOOGLE_API_KEY)
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("Notes_faiss_index")
        return vector_store
    except Exception as e:
        st.error(f"Error creating vector store: {str(e)}")
        return None

def clear_vector_store():
    try:
        os.remove("Notes_faiss_index")
    except FileNotFoundError:
        pass
    except Exception as e:
        st.error(f"Error clearing vector store: {str(e)}")

def get_conversational_chain():
    """Creates the question-answering chain."""
    try:
        prompt_template = """
        As a skilled notetaking and teaching expert, your task is to provide comprehensive answers based on the provided context. Your goal is to assist the user in preparing for an exam by creating a structured learning path with necessary concepts. Ensure that your answers cover all pertinent details and address the user's questions effectively.\n\n
        Context:\n {context}?\n
        Question:\n {question}\n
        
        Answer:
        """
        model = ChatGoogleGenerativeAI(model='gemini-pro', temperature=0.3, google_api_key=GOOGLE_API_KEY)
        prompt = PromptTemplate(template=prompt_template, input_variables=["context", 'question'])
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
    except Exception as e:
        st.error(f"Error loading conversational chain: {str(e)}")
        return None

    
def user_input(user_question):
    """Processes user question and retrieves answer from PDFs."""
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=GOOGLE_API_KEY)
        new_db = FAISS.load_local("faiss_index", embeddings=embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)

        chain = get_conversational_chain()
        if chain:
            response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
            st.write("Reply: ", response["output_text"])
        else:
            st.error("Error processing user input.")
    except Exception as e:
        st.error(f"Error processing user input: {str(e)}")

def main():
    try:
        st.set_page_config(page_title="MyThoughts.AI", layout='wide', initial_sidebar_state='auto')
        st.title("Welcome to [MyThoughts.AI](https://github.com/vibhansh/Langchain_Gemini_Test_App)")
        st.write("Interact with your data",)
        
        # Empty placeholder for user response
        user_response = st.empty()

        # Sidebar for PDF upload and processing
        with st.sidebar:
            st.title("Menu:")
            pdf_docs = st.file_uploader("Upload your PDF docs and click Submit", accept_multiple_files=True)
            url = st.text_input("Enter YouTube video URL")
            slide_data = st.file_uploader("Upload Slide data", type=['pptx', 'ppt'])
            text_data = st.text_input("Paste your text data")

            if st.button("Submit & Process"):
                with st.spinner("Processing..."):
                    if pdf_docs:
                        for pdf_file in pdf_docs:
                            raw_text = get_pdf_text(pdf_file.getvalue())
                            text_chunks = get_text_chunks(raw_text)
                            get_vector_store(text_chunks)
                    if url:
                        raw_text = get_transcript_from_captions(url)
                        get_vector_store([raw_text])  # Wrap in a list for consistency
                    if slide_data:
                        raw_text = process_google_slides(slide_data)
                        get_vector_store([raw_text])
                    if text_data:
                        get_vector_store([text_data])
                    st.success("Done")

        # Text input for user question
        user_question = st.text_input("Ask a Question about the uploaded data")

        # Update user response based on user input
        if user_question:
            user_input(user_question)
    except Exception as e:
        st.error(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    main()
